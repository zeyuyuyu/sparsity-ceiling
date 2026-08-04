#!/usr/bin/env python
"""patch_deck3.py — shape-law/streaming narrowing (2026-08-04).

Anchored discipline (same as patch_deck2.py): every edit asserts the shape's
CURRENT text before touching it; any drift aborts the whole run with no
partial application. Formatting preserved by writing into the first run.

Edits (numbers from ledger entries d47511b / c18c724 / 0924aad, none re-derived):
  E1 slide 7 (07 - SSM RESULTS), shape id 50: band-05 title -> the unqualified verdict.
  E2 slide 7, shape id 51: band-05 body gains the shape law + streaming negative.
  E3 slide 8 (08 - PROPOSED DIRECTIONS), shape id 9: direction-01 body no longer
     implies the streaming shape is untested (it is refuted, budget-robust).
"""
import sys
from pptx import Presentation

PPTX = "beyond_attention_paradigms.pptx"

def norm(s):
    return (s.replace("’", "'").replace("‘", "'")
             .replace("“", '"').replace("”", '"')
             .replace("–", "-").replace("—", "-")
             .replace("×", "x").replace("≤", "<=")
             .replace("…", "...").replace("≪", "<<")
             .replace("→", "->")).strip()

EDITS = [
    # (slide_idx, shape_id, expected_norm_text, new_text)
    (7, 50,
     "The energy is not where the mechanisms are",
     "No quality-matched pJ win at any shape tested"),
    (7, 51,
     ("The readout carries 75% of analog's pJ/token, the recurrence 8%; even a "
      "free recurrence costs 1.74x. Gating the readout recovers the energy but "
      "costs ~+1.0 bpc (char-LM) and all of copy's margin, 3 seeds. Converters: 0.00-0.33%."),
     ("LM shape: readout = 75% of analog’s pJ/token, recurrence 8%; gating it "
      "costs ~+1.0 bpc and all of copy’s margin (3 seeds). Streaming shape "
      "(V/H=0.04): recurrence dominates, 3.3× recoverable — but −10…−19 pts, "
      "budget-robust. Converters ≤0.58%.")),
    (8, 9,
     ("Finished in simulation, including the energy datapath: the readout's 75% "
      "share resisted every gating mechanism tested. What is unfinished is "
      "physical - every pJ figure is a 45nm proxy. So 02 first, not more simulation."),
     ("Done in simulation at both shapes: LM — the readout’s 75% share resisted "
      "every gating mechanism; streaming — 3.3× is recoverable but costs 10–19 "
      "pts, budget-robust. What’s unfinished is physical: every pJ is a 45nm "
      "proxy. So 02 first, not more simulation.")),
]

prs = Presentation(PPTX)

# Pass 1: verify every anchor before changing anything.
targets = []
for slide_idx, shape_id, expected, new in EDITS:
    shape = next((sh for sh in prs.slides[slide_idx].shapes if sh.shape_id == shape_id), None)
    if shape is None or not shape.has_text_frame:
        sys.exit(f"ABORT: slide {slide_idx} shape id {shape_id} missing or not text")
    cur = norm(shape.text_frame.text)
    if cur != norm(expected):
        sys.exit(f"ABORT: slide {slide_idx} shape {shape_id} drifted.\n  have: {cur!r}\n  want: {norm(expected)!r}")
    targets.append((shape, new))

# Pass 2: apply, preserving run formatting (all text into run 0, blank the rest).
for shape, new in targets:
    tf = shape.text_frame
    first = True
    for para in tf.paragraphs:
        for run in para.runs:
            if first:
                run.text = new
                first = False
            else:
                run.text = ""
    if first:
        sys.exit("ABORT: target shape had no runs")

prs.save(PPTX)
print("OK: 3 edits applied ->", PPTX)
