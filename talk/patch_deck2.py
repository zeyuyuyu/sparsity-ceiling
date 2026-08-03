#!/usr/bin/env python
"""Patch the talk deck with the energy-datapath / stranded-readout result.

Adds a fifth block to the SSM RESULTS slide and de-stales the directions
slide's "converters are unpriced" line (they are now priced under the proxy).

Every number is copied from published ledger entries (c8b596f/a9ff715,
73c99e6, 60d627c, db7c13b) and paper2 sec 8.4, never re-derived.  The script
matches every target by shape id AND asserts its current text/geometry, so a
drifted deck aborts instead of half-applying.

CLAIMS DISCIPLINE: the readout result is a NEGATIVE (no pJ win); energy is a
45nm proxy, not silicon; converters are priced under that proxy, not measured.
"""
import copy
import sys

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches

PPTX = "/work/zeyuwang/sparsity-ceiling/talk/beyond_attention_paradigms.pptx"
EPS = Inches(0.02)

prs = Presentation(PPTX)
if len(prs.slides) != 10:
    sys.exit("ABORT: expected 10 slides, found %d" % len(prs.slides))

res = prs.slides[7]
dirs = prs.slides[8]


def norm(s):
    return " ".join(s.split())


def by_id(slide, sid):
    hits = [s for s in slide.shapes if s.shape_id == sid]
    if len(hits) != 1:
        sys.exit("ABORT: shape id %d matched %d shapes" % (sid, len(hits)))
    return hits[0]


def check(shape, text, left, top):
    got = norm(shape.text_frame.text) if shape.has_text_frame else ""
    if not got.startswith(norm(text)):
        sys.exit("ABORT: id %d text %r != expected %r" % (shape.shape_id, got[:60], text[:60]))
    if abs(shape.left - Inches(left)) > EPS or abs(shape.top - Inches(top)) > EPS:
        sys.exit(
            "ABORT: id %d at (%.2f,%.2f) != expected (%.2f,%.2f)"
            % (shape.shape_id, Emu(shape.left).inches, Emu(shape.top).inches, left, top)
        )
    return shape


def set_text(shape, new):
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        sys.exit("ABORT: shape %d has no runs to inherit formatting from" % shape.shape_id)
    p0.runs[0].text = new
    for r in p0.runs[1:]:
        r._r.getparent().remove(r._r)
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)


# --------------------------------------------------- verify the results slide is the one we think
check(by_id(res, 2), "07 · SSM RESULTS", 0.60, 0.50)
check(by_id(res, 5), "Direction One, run: what the analog-state SSM actually did", 0.60, 0.90)

ROW1 = [(6, "", 0.60, 1.80), (7, "01", 0.90, 2.02),
        (8, "Analog state buys energy, not quality", 1.65, 2.04),
        (9, "Param-matched 4-way SSM", 1.65, 2.62),
        (10, "", 6.85, 1.80), (11, "02", 7.15, 2.02),
        (12, "The variant ranking inverts by workload", 7.90, 2.04),
        (13, "At matched communication rate", 7.90, 2.62)]
ROW2 = [(14, "", 0.60, 3.85), (15, "03", 0.90, 4.07),
        (16, "Degrade what the task ignores", 1.65, 4.09),
        (17, "One principle covers both", 1.65, 4.67),
        (18, "", 6.85, 3.85), (19, "04", 7.15, 4.07),
        (20, "My own bound did not survive its test", 7.90, 4.09),
        (21, "Shrink-H, 3 seeds per width", 7.90, 4.67)]

for sid, txt, l, t in ROW1 + ROW2:
    check(by_id(res, sid), txt, l, t)

# --------------------------------------------------- clone block 03 into a slim full-width band 05
# The four existing blocks stay put; the band occupies the empty strip between
# the second row (ends 5.70) and the footer (7.00), with title and body side by
# side so it needs 1.00in of height instead of the blocks' 1.85in.
src_ids = [14, 15, 16, 17]  # rect, number, title, body
next_id = max(s.shape_id for sl in prs.slides for s in sl.shapes) + 1
clones = []
for sid in src_ids:
    el = copy.deepcopy(by_id(res, sid)._element)
    res.shapes._spTree.append(el)
    clones.append(el)

# the clones carry duplicate ids until renamed, so resolve them positionally, not by id
new_shapes = list(res.shapes)[-4:]
for shp in new_shapes:
    cnvpr = shp._element.find(qn("p:nvSpPr")).find(qn("p:cNvPr"))
    cnvpr.set("id", str(next_id))
    cnvpr.set("name", "energy-block-%d" % next_id)
    next_id += 1
new_rect, new_num, new_title, new_body = new_shapes

TOP = Inches(5.85)
new_rect.left, new_rect.top = Inches(0.60), TOP
new_rect.width, new_rect.height = Inches(12.20), Inches(1.00)
new_num.left, new_num.top = Inches(0.90), TOP + Inches(0.18)
new_title.left, new_title.top = Inches(1.65), TOP + Inches(0.14)
new_title.width, new_title.height = Inches(4.25), Inches(0.75)
new_body.left, new_body.top = Inches(6.05), TOP + Inches(0.14)
new_body.width, new_body.height = Inches(6.60), Inches(0.78)

set_text(new_num, "05")
set_text(new_title, "The energy is not where the mechanisms are")
set_text(
    new_body,
    "The readout carries 75% of analog's pJ/token, the recurrence 8%; even a free recurrence costs "
    "1.74×. Gating the readout recovers the energy but costs ~+1.0 bpc (char-LM) and all of copy's "
    "margin, 3 seeds. Converters: 0.00–0.33%.",
)

# --------------------------------------------------- directions slide: converters are no longer unpriced
check(by_id(dirs, 2), "08 · PROPOSED DIRECTIONS", 0.60, 0.50)
old = check(
    by_id(dirs, 9),
    "Finished in simulation (two tasks, three seeds); what is unfinished is physical",
    1.65,
    2.62,
)
set_text(
    old,
    "Finished in simulation, including the energy datapath: the readout's 75% share resisted every "
    "gating mechanism tested. What is unfinished is physical — every pJ figure is a 45nm proxy. So "
    "02 first, not more simulation.",
)

prs.save(PPTX)
print("OK: patched, %d slides, %d shapes on results slide" % (len(prs.slides), len(list(res.shapes))))
