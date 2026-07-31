#!/usr/bin/env python
"""Patch talk/beyond_attention_paradigms.pptx to the post-retraction SSM framing.

Every number here is copied from published ledger entries / imam_ssm_memo_v3.md,
never re-derived.  The script aborts if any text anchor is missing, so it can
never half-apply.

CLAIMS DISCIPLINE: never "analog beats digital" (retracted), never "bound
confirmed", floor swing is 12x not 190x, copy margins are vs the
noise-regularized reference, energy is a 45nm proxy not silicon.
"""
import copy
import sys

from pptx import Presentation
from pptx.oxml.ns import qn

PPTX = "/work/zeyuwang/sparsity-ceiling/talk/beyond_attention_paradigms.pptx"

prs = Presentation(PPTX)
if len(prs.slides) != 9:
    sys.exit("ABORT: expected 9 slides, found %d" % len(prs.slides))


def shapes_with_text(slide):
    return [s for s in slide.shapes if s.has_text_frame]


def norm(s):
    """Compare ignoring the exotic spaces the deck uses (e.g. U+202F in "Loihi 2")."""
    return " ".join(s.split())


def find(slide, text, where=""):
    hits = [s for s in shapes_with_text(slide) if norm(s.text_frame.text) == norm(text)]
    if len(hits) != 1:
        sys.exit("ABORT: anchor %r matched %d shapes %s" % (text[:60], len(hits), where))
    return hits[0]


def set_text(shape, new):
    """Replace a shape's text, keeping the first run's formatting."""
    tf = shape.text_frame
    paras = tf.paragraphs
    p0 = paras[0]
    if not p0.runs:
        sys.exit("ABORT: shape has no runs to inherit formatting from")
    r0 = p0.runs[0]
    r0.text = new
    for r in p0.runs[1:]:
        r._r.getparent().remove(r._r)
    for p in paras[1:]:
        p._p.getparent().remove(p._p)


# ---------------------------------------------------------------- slide 5: bound annotation
s5 = prs.slides[5]
set_text(
    find(s5, "memory in membranes/synapses, not a KV cache → no memory wall, no firing floor", "s5"),
    "memory in membranes/synapses, not a KV cache → no memory wall; the firing floor is a "
    "spiking-state property (the bound itself: theory, scope limit tested)",
)

# ---------------------------------------------------------------- new results slide (clone of the directions slide)
src = prs.slides[7]
new = prs.slides.add_slide(src.slide_layout)
for shp in list(new.shapes):
    shp._element.getparent().remove(shp._element)
for shp in src.shapes:
    new.shapes._spTree.append(copy.deepcopy(shp._element))

# the dark background lives on the slide (p:bg), not in the shape tree
src_csld = src._element.find(qn("p:cSld"))
src_bg = src_csld.find(qn("p:bg"))
if src_bg is None:
    sys.exit("ABORT: source slide has no p:bg to clone")
new._element.find(qn("p:cSld")).insert(0, copy.deepcopy(src_bg))

sld_lst = prs.slides._sldIdLst
ids = list(sld_lst)
sld_lst.remove(ids[-1])
sld_lst.insert(7, ids[-1])

RESULTS = [
    (
        "Escape the firing floor with analog state",
        "Analog state buys energy, not quality",
        "Param-matched 4-way SSM, char-LM, 3 seeds: analog state cuts proxy energy 1.60× (446k vs "
        "712k pJ/token) for +0.160 bpc. “Beats digital” is retracted — the same noise on the digital "
        "baseline is worth −0.309 bpc.",
    ),
    (
        "Predictive coding on neuromorphic silicon",
        "The variant ranking inverts by workload",
        "At matched communication rate the order inverts: char-LM analog ≪ spiking-state ≪ "
        "output-spiking; precise recall (copy) the reverse — margin kept 0.87 / 0.50 / 0.074. A 1-bit "
        "state is at exactly chance, and cheapest.",
    ),
    (
        "Quantify the floor–vs–wall trade-off",
        "Degrade what the task ignores",
        "One principle covers both: a degradation is cheap only where the loss does not depend on "
        "it. The same 6-bit state quantizer costs −0.002 bpc on char-LM and +0.37 on copy (post-hoc, "
        "not pre-registered).",
    ),
    (
        "Closed-loop world-model at the edge",
        "My own bound did not survive its test",
        "Shrink-H, 3 seeds per width: the predicted floor rises 12× while measured state activity "
        "falls 0.496→0.250, and accuracy improves as capacity shrinks. Theory with a tested scope "
        "limit, not a validated prediction.",
    ),
]

set_text(find(new, "07 · PROPOSED DIRECTIONS", "new"), "07 · SSM RESULTS")
set_text(find(new, "08", "new"), "08")
set_text(
    find(new, "What I would work on", "new"),
    "Direction One, run: what the analog-state SSM actually did",
)
for old_title, new_title, body in RESULTS:
    set_text(find(new, old_title, "new"), new_title)

OLD_BODIES = [
    "Build an SSM whose recurrent state lives in sub-threshold analog dynamics, not spikes — test "
    "whether it breaks our bound's assumption and sparsifies like attention did.",
    "Implement an error-propagating hierarchy with local (EqProp) learning; measure real energy on "
    "Loihi 2 / SpiNNaker2 — turning the 45nm proxy into a measurement.",
    "Formalize the recurrence(firing floor) ↔ attention(memory wall) dichotomy as a conservation "
    "law; find the Pareto frontier and where analog-state models sit on it.",
    "A predictive world-model + planning loop for a low-power embodied agent — where neuromorphic's "
    "real-time, event-driven physics is decisive.",
]
for old_body, (_, _, body) in zip(OLD_BODIES, RESULTS):
    set_text(find(new, old_body, "new"), body)

# ---------------------------------------------------------------- directions slide: renumber + fold direction 01
d = prs.slides[8]
set_text(find(d, "07 · PROPOSED DIRECTIONS", "dir"), "08 · PROPOSED DIRECTIONS")
set_text(find(d, "08", "dir"), "09")
set_text(find(d, "What I would work on", "dir"), "What I would work on next")
set_text(
    find(d, "Escape the firing floor with analog state", "dir"),
    "Analog state — done in simulation, open in silicon",
)
set_text(
    find(d, OLD_BODIES[0], "dir"),
    "Finished in simulation (two tasks, three seeds); what is unfinished is physical — every pJ "
    "figure is a 45nm proxy and the analog storage element and converter are unpriced. So 02 "
    "first, not more simulation.",
)

prs.save(PPTX)
print("OK: patched, %d slides" % len(prs.slides))
